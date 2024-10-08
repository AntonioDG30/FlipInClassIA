import uuid

import openai
from flask import Flask, render_template, request, redirect, url_for, session, flash, jsonify
from PIL import Image, ImageDraw, ImageFont
from datetime import datetime
from pptx import Presentation
import pdfplumber
import mysql.connector
import bcrypt
import os
import secrets
import random
import re
import nltk
import numpy as np

nltk.download('punkt')  # Scarica solo il pacchetto necessario

# Inizializza l'app Flask
app = Flask(__name__)

# Genera una chiave segreta casuale per la sessione dell'app
app.secret_key = secrets.token_bytes(32)

# Percorso di upload
UPLOAD_FOLDER = 'static/fileCaricati/'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# Percorsi per le immagini
IMAGES_PATH = "static/dashboardStaticFile/img/corsi"
os.makedirs(IMAGES_PATH, exist_ok=True)
USER_IMAGES_PATH = "static/dashboardStaticFile/img/utenti"
os.makedirs(USER_IMAGES_PATH, exist_ok=True)

# Configura la chiave API di OpenAI
openai.api_key = os.getenv('OPENAI_API_KEY')


# Funzione migliorata per generare il feedback personalizzato per lo studente
def genera_feedback(domande, risposte):
    try:
        # Richiesta per generare le 5 macro-aree più importanti
        prompt_macro = (
            f"Per favore, agisci come se fossi un professore universitario e restituisci un feedback in italiano agli studenti "
            f"per aiutarli a comprendere cosa devono migliorare e come farlo. Il feedback deve essere una frase da massimo 2/3 righe. "
            f"Rispondimi restituendomi solo il feedback. Il feedback deve essere mirato a far capire cosa lo studente non ha saputo "
            f"rispondere, analizzando magari l'argomento errato ecc. Considera ha svolto il questionario due volte, una ad inizio lezione ed una alla fine. "
            f"Realizza il feedback basandoti sulle seguenti domande: {domande} al cui lo studente ha dato le seguenti risposte: {risposte}")

        response_macro = openai.ChatCompletion.create(
            model="gpt-4",
            messages=[
                {"role": "system",
                 "content": "Sei un professore universitario che genera feedback personalizzati ai propri studenti dalle risposte"
                            " di un questionario."},
                {"role": "user", "content": prompt_macro}
            ],
            max_tokens=300,  # Limita i token di output
            temperature=0.5,
        )

        # Estrai il contenuto del feedback dalla risposta dell'API
        feedback = response_macro.choices[0].message['content']

    except Exception as e:
        print(f"Errore durante la generazione: {e}")
        return ""

    return feedback  # Ritorna il feedback estratto


@app.route('/ottieniFeedback')
def ottieniFeedback():
    # Controlla se l'utente è loggato
    if 'user_id' in session:
        user_id = session['user_id']  # Recupera l'ID utente dalla sessione
        user_type = session['user_type']  # Recupera il tipo di utente (studente o docente)

        # Recupera il corso_id e questionario_id passati come parametri dalla query string
        questionario_id = request.args.get('questionario_id')
        corso_id = request.args.get('corso_id')

        # Connessione al database
        conn = get_db_connection()
        if conn is None:
            # In caso di errore nella connessione, restituisce un messaggio di errore come JSON
            return jsonify({"error": "Errore di connessione al database"}), 500

        cursor = conn.cursor(dictionary=True)

        # Ottenere le domande e le opzioni della lezione
        query_domande = """
                    SELECT d.domanda_id, d.testo_domanda, d.corretta_opzione_id
                    FROM Domanda d
                    WHERE d.questionario_id = %s
                """
        cursor.execute(query_domande, (questionario_id,))
        domande = cursor.fetchall()

        # Aggiungere le opzioni per ogni domanda
        for domanda in domande:
            query_opzioni = """
                        SELECT o.testo_opzione, o.opzione_id
                        FROM Opzione o
                        WHERE o.domanda_id = %s
                    """
            cursor.execute(query_opzioni, (domanda['domanda_id'],))
            domanda['opzioni'] = cursor.fetchall()

        query_risposte = """
                            SELECT r.domanda_id, r.opzione_scelta_id
                            FROM risposta_studente r
                            JOIN domanda d ON d.domanda_id = r.domanda_id
                            WHERE r.studente_id = %s and d.questionario_id = %s
                        """
        cursor.execute(query_risposte, (user_id, questionario_id,))
        risposte = cursor.fetchall()

        feedback = genera_feedback(domande, risposte)

        query_statistiche_id = """
                                    SELECT s.statistiche_studente_id
                                    FROM statistiche_studente s
                                    WHERE s.studente_id = %s and s.questionario_id = %s
                                    ORDER BY s.statistiche_studente_id DESC LIMIT 1
                                """
        cursor.execute(query_statistiche_id, (user_id, questionario_id,))
        statistiche_studente_id = cursor.fetchone()['statistiche_studente_id']

        salva_feedback = """
                              UPDATE statistiche_studente s
                              SET s.feedback_suggerito = %s
                              WHERE s.statistiche_studente_id = %s                                
                         """
        cursor.execute(salva_feedback, (feedback, statistiche_studente_id,))

        # Chiude il cursore e la connessione
        cursor.close()
        conn.commit()
        conn.close()

        # Restituisce la pagina delle statistiche ('ottieniStatistiche.html')
        return redirect(url_for('ottieniStatistiche', corso_id=corso_id))

    else:
        # Se l'utente non è loggato, mostra un messaggio di errore e reindirizza alla pagina di login
        flash('Devi essere loggato per accedere alla dashboard.')
    return redirect(url_for('login'))


# Rotta per visualizzare la lista dei professori che partecipano a un corso
@app.route('/ottieniStatistiche')
def ottieniStatistiche():
    # Controlla se l'utente è loggato
    if 'user_id' in session:
        user_id = session['user_id']  # Recupera l'ID utente dalla sessione
        user_type = session['user_type']  # Recupera il tipo di utente (studente o docente)

        # Recupera il corso_id passato come parametro dalla query string
        corso_id = request.args.get('corso_id')

        # Connessione al database
        conn = get_db_connection()
        if conn is None:
            # In caso di errore nella connessione, restituisce un messaggio di errore come JSON
            return jsonify({"error": "Errore di connessione al database"}), 500

        cursor = conn.cursor(dictionary=True)

        # Ottieni tutte le statistiche
        statistiche_questionario = []
        statistiche_studente = []

        if user_type == 'docente':
            # Ottieni tutte le statistiche dei questionari per il corso
            cursor.execute("""
                    SELECT l.data AS data, sq.punteggio_medio, sq.percentuale_successo, 
                           sq.numero_domande, sq.risposte_corrette, sq.risposte_errate, sq.varianza_punteggio
                    FROM questionario q
                    JOIN statistiche_questionario sq ON q.questionario_id = sq.questionario_id
                    JOIN lezione l ON q.lezione_id = l.lezione_id
                    WHERE l.corso_id = %s
                    ORDER BY l.data
                """, (corso_id,))
            statistiche_questionario = cursor.fetchall()

        elif user_type == 'studente':
            # Ottieni tutte le statistiche per lo studente
            cursor.execute("""
                    SELECT l.data AS data, ss.punteggio, ss.risposte_corrette, ss.risposte_errate, ss.feedback_suggerito
                    FROM questionario q
                    JOIN statistiche_studente ss ON q.questionario_id = ss.questionario_id
                    JOIN lezione l ON q.lezione_id = l.lezione_id
                    WHERE ss.studente_id = %s AND l.corso_id = %s
                    ORDER BY l.data
                """, (user_id, corso_id))
            statistiche_studente = cursor.fetchall()

        cursor.close()
        conn.close()

        docente_presidente = ottieniProfessorePresidente(corso_id)  # Ottiene il docente responsabile del corso
        nome_corso = ottieniNomeCorso(corso_id)  # Ottiene il nome del corso

        # Chiude il cursore e la connessione
        cursor.close()
        conn.close()

        # Restituisce la pagina con l'elenco dei professori partecipanti ('professori_partecipanti.html')
        return render_template('statistiche.html', user_id=user_id, user_type=user_type,
                               corso_id=corso_id, statistiche_questionario=statistiche_questionario,
                               statistiche_studente=statistiche_studente,
                               docente_presidente=docente_presidente, nome_corso=nome_corso)

    else:
        # Se l'utente non è loggato, mostra un messaggio di errore e reindirizza alla pagina di login
        flash('Devi essere loggato per accedere alla dashboard.')
    return redirect(url_for('login'))


@app.route('/invia_questionario', methods=['POST'])
def invia_questionario():
    data = request.json
    lezione_id = data.get('lezione_id')
    studente_id = data.get('studente_id')
    questionario_id = data.get('questionario_id')
    risposte = data.get('risposte', [])

    try:
        # Connessione al database
        conn = get_db_connection()
        cursor = conn.cursor()

        # Aggiorna la tabella "Risposta_Studente"
        for risposta in risposte:
            domanda_id = risposta['domanda_id']
            opzione_scelta_id = risposta['opzione_scelta_id']

            # Inserisci nella tabella "Risposta_Studente"
            cursor.execute('''
                INSERT INTO risposta_studente (studente_id, domanda_id, opzione_scelta_id)
                VALUES (%s, %s, %s)
            ''', (studente_id, int(domanda_id), int(opzione_scelta_id) if opzione_scelta_id is not None else None))

        # Calcola le statistiche dello studente
        num_risposte_corrette = sum(1 for r in risposte if r['risposta'] == r['corretta'])
        num_risposte_errate = len(risposte) - num_risposte_corrette
        punteggio = (num_risposte_corrette / len(risposte)) * 100

        # Aggiorna "Statistiche_Studente"
        cursor.execute('''
            INSERT INTO statistiche_studente (studente_id, questionario_id, risposte_corrette, risposte_errate, punteggio)
            VALUES (%s, %s, %s, %s, %s)
            ON DUPLICATE KEY UPDATE
            risposte_corrette = VALUES(risposte_corrette), 
            risposte_errate = VALUES(risposte_errate), 
            punteggio = VALUES(punteggio)
        ''', (studente_id, int(questionario_id), num_risposte_corrette, num_risposte_errate, punteggio))

        # Ottieni tutte le statistiche degli studenti per calcolare i valori aggregati
        cursor.execute('''
            SELECT punteggio
            FROM statistiche_studente
            WHERE questionario_id = %s
        ''', (int(questionario_id),))
        punteggi_studenti = cursor.fetchall()
        punteggi_studenti = [p[0] for p in punteggi_studenti]  # Convertiamo i punteggi in una lista di float

        # Calcola `percentuale_successo`
        num_studenti_successo = sum(1 for p in punteggi_studenti if p >= 50)
        percentuale_successo = (num_studenti_successo / len(punteggi_studenti)) * 100 if punteggi_studenti else 0

        # Calcola `varianza_punteggio`
        varianza_punteggio = np.var(punteggi_studenti) if punteggi_studenti else 0

        # Ottieni le statistiche attuali per aggiornare "Statistiche_Questionario"
        cursor.execute('''
            SELECT numero_domande, risposte_corrette, risposte_errate, punteggio_medio
            FROM statistiche_questionario
            WHERE questionario_id = %s
        ''', (int(questionario_id),))
        result = cursor.fetchone()

        if result:
            numero_domande, risposte_corrette, risposte_errate, punteggio_medio = result

            # Calcola i nuovi valori
            nuovo_numero_domande = numero_domande + len(risposte)
            nuovo_risposte_corrette = risposte_corrette + num_risposte_corrette
            nuovo_risposte_errate = risposte_errate + num_risposte_errate
            nuovo_punteggio_medio = (nuovo_risposte_corrette * 100) / nuovo_numero_domande

            # Aggiorna "Statistiche_Questionario"
            cursor.execute('''
                UPDATE statistiche_questionario 
                SET numero_domande = %s, 
                    risposte_corrette = %s, 
                    risposte_errate = %s,
                    punteggio_medio = %s,
                    percentuale_successo = %s,
                    varianza_punteggio = %s
                WHERE questionario_id = %s
            ''', (nuovo_numero_domande, nuovo_risposte_corrette, nuovo_risposte_errate, nuovo_punteggio_medio,
                  percentuale_successo, varianza_punteggio, int(questionario_id)))
        else:
            # Inserisci un nuovo record se non esiste ancora
            cursor.execute('''
                INSERT INTO statistiche_questionario (questionario_id, numero_domande, risposte_corrette, risposte_errate, punteggio_medio, percentuale_successo, varianza_punteggio)
                VALUES (%s, %s, %s, %s, %s, %s, %s)
            ''', (int(questionario_id), len(risposte), num_risposte_corrette, num_risposte_errate, punteggio,
                  percentuale_successo, varianza_punteggio))

        # Commit delle modifiche
        conn.commit()
        cursor.close()
        conn.close()

        return jsonify({'success': True})

    except Exception as e:
        print(f'Errore durante l\'aggiornamento del database: {e}')
        return jsonify({'success': False, 'error': str(e)}), 500


# Endpoint API per ottenere lo stato della lezione
@app.route('/aggiorna_dati_studenti', methods=['POST'])
def aggiorna_dati_studenti():
    lezione_id = request.json.get('lezione_id')
    questionario_id = request.json.get('questionario_id')

    # Connessione al database
    conn = get_db_connection()
    if conn is None:
        return jsonify({"error": "Errore di connessione al database"}), 500

    cursor = conn.cursor(dictionary=True)

    # 1. Ottenere l'elenco degli studenti presenti alla lezione
    query = "SELECT statolezione FROM lezione WHERE lezione_id = %s"
    cursor.execute(query, (lezione_id,))
    statolezione = cursor.fetchone()

    # 2. Ottenere ora inizio questionario
    query = "SELECT avvio_primaFase, avvio_secondaFase FROM questionario WHERE lezione_id = %s AND questionario_id = %s"
    cursor.execute(query, (lezione_id, questionario_id))
    avvioquestionario = cursor.fetchone()

    cursor.close()
    conn.close()

    return jsonify({
        'statolezione': statolezione['statolezione'],
        'avvio_primaFase': int(avvioquestionario['avvio_primaFase'].timestamp()) if avvioquestionario[
            'avvio_primaFase'] else None,
        'avvio_secondaFase': int(avvioquestionario['avvio_secondaFase'].timestamp()) if avvioquestionario[
            'avvio_secondaFase'] else None
    })


@app.route('/ottieni_questionario', methods=['POST'])
def ottieni_questionario():
    lezione_id = request.json.get('lezione_id')

    if lezione_id is None:
        return jsonify({"error": "ID lezione non fornito"}), 400

    # Connessione al database
    conn = get_db_connection()
    if conn is None:
        return jsonify({"error": "Errore di connessione al database"}), 500

    cursor = conn.cursor(dictionary=True)

    # Ottenere le domande e le opzioni della lezione
    query_domande = """
            SELECT d.domanda_id, d.testo_domanda, d.corretta_opzione_id
            FROM Domanda d
            JOIN Questionario q ON d.questionario_id = q.questionario_id
            WHERE q.lezione_id = %s
        """
    cursor.execute(query_domande, (lezione_id,))
    domande = cursor.fetchall()

    # Aggiungere le opzioni per ogni domanda
    for domanda in domande:
        query_opzioni = """
                SELECT o.testo_opzione, o.opzione_id
                FROM Opzione o
                WHERE o.domanda_id = %s
            """
        cursor.execute(query_opzioni, (domanda['domanda_id'],))
        domanda['opzioni'] = cursor.fetchall()

    # Chiude il cursore e la connessione
    cursor.close()
    conn.close()

    return jsonify({
        'domande': domande
    })


# Funzione per estrarre testo da PDF
def estrai_testo_da_pdf(filepath):
    testo = ""
    with pdfplumber.open(filepath) as pdf:
        for pagina in pdf.pages:
            testo += pagina.extract_text() + "\n"
    return testo


# Funzione per estrarre testo da PPTX
def estrai_testo_da_pptx(filepath):
    testo = ""
    prs = Presentation(filepath)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                testo += shape.text + "\n"
    return testo


# Funzione migliorata per pulire e segmentare il testo
def pulisci_e_segmenta_testo(testo):
    # Rimuovo caratteri non ASCII e spazi in eccesso
    testo_pulito = re.sub(r'[^\x00-\x7F]+', ' ', testo)
    testo_pulito = re.sub(r'\s+', ' ', testo_pulito).strip()

    # Segmentazione del testo in frasi
    frasi = nltk.sent_tokenize(testo_pulito, language="italian")

    # Segmenta il testo in parti più piccole
    segmenti = []
    segmento_corrente = ""
    token_count = 0

    for frase in frasi:
        num_tokens = len(frase.split())
        if token_count + num_tokens <= 1500:  # Limita ogni segmento a circa 1500 parole
            segmento_corrente += frase + " "
            token_count += num_tokens
        else:
            segmenti.append(segmento_corrente.strip())
            segmento_corrente = frase + " "
            token_count = num_tokens

    # Aggiunge l'ultimo segmento
    if segmento_corrente:
        segmenti.append(segmento_corrente.strip())

    return segmenti


def genera_macro_aree_e_questionario(testo):
    # Segmenta il testo in parti più piccole
    segmenti = pulisci_e_segmenta_testo(testo)

    macro_aree = []
    questionario = []
    domande_generate = 0
    domande_gia_generate = []  # Registro delle domande già generate

    # Unisci solo i primi 2 segmenti e riduci il testo totale per evitare di superare il limite di token
    testo_unificato = " ".join(segmenti[:2])[:2000]

    try:
        # Richiesta per generare le 5 macro-aree più importanti
        prompt_macro = (
            f"Per favore, agisci come se fossi un professore universitario e ricava le 5 macro-aree/argomenti in italiano più importanti, "
            f"Ogni macro-area dovrebbe essere nel formato 'Nome Argomento: Descrizione dell'argomento' e non numerare le macro-aree, "
            f"della lezione in italiano riassunta nel seguente testo: {testo_unificato}")

        response_macro = openai.ChatCompletion.create(
            model="gpt-4",
            messages=[
                {"role": "system",
                 "content": "Sei un professore universitario che genera macro-aree da un testo di lezione."},
                {"role": "user", "content": prompt_macro}
            ],
            max_tokens=300,
            temperature=0.5,
        )

        # Supponiamo che l'output sia una lista di argomenti nel formato 'Nome Argomento: Descrizione'
        macro_lines = response_macro['choices'][0]['message']['content'].strip().split('\n')
        for line in macro_lines:
            if ':' in line:
                nome_argomento, descrizione_argomento = line.split(':', 1)
                nome_argomento = nome_argomento.strip()
                descrizione_argomento = descrizione_argomento.strip()
                if nome_argomento and descrizione_argomento:
                    macro_aree.append((nome_argomento, descrizione_argomento))

        # Limita il numero di domande da generare tra 10 e 15
        num_domande = random.randint(10, 15)

        # Suddividi ulteriormente il testo in segmenti più piccoli per la generazione del questionario
        for segmento in segmenti:
            if domande_generate >= num_domande:
                break

            segmento_limitato = segmento[:1000]

            while domande_generate < num_domande:
                num_domande_richieste = min(5, num_domande - domande_generate)

                # Crea una lista di domande già generate per informare l'AI di evitare ripetizioni
                domande_testo = "\n".join([domanda['testo_domanda'] for domanda in domande_gia_generate])

                # Richiesta per generare un blocco di domande
                prompt_questionario = (
                    f"Basandoti sul testo seguente, crea un questionario a risposta multipla con 4 opzioni di cui una sola è corretta. "
                    f"Crea {num_domande_richieste} domande che variano su tutti gli argomenti del testo. Non numerare le domande. "
                    f"Evita di ripetere i seguenti argomenti: {domande_testo}. "
                    f"Assicurati che ogni domanda abbia esattamente 4 opzioni e una risposta corretta indicata chiaramente nel formato 'Risposta corretta: X)'. "
                    f"Formato delle domande: [Domanda], A) [Opzione 1], B) [Opzione 2], C) [Opzione 3], D) [Opzione 4], Risposta corretta: X). "
                    f"Ecco il testo: {segmento_limitato}")

                response_questionario = openai.ChatCompletion.create(
                    model="gpt-4",
                    messages=[
                        {"role": "system",
                         "content": "Sei un professore universitario che crea questionari a risposta multipla basati su un testo."},
                        {"role": "user", "content": prompt_questionario}
                    ],
                    max_tokens=500,
                    temperature=0.5,
                )

                # Analizza la risposta per estrarre le domande e le opzioni
                question_lines = response_questionario['choices'][0]['message']['content'].strip().split('\n')
                current_question = None
                corretta_opzione = None

                for line in question_lines:
                    line = line.strip()

                    # Identifica una nuova domanda se la linea non inizia con un'opzione o con 'Risposta corretta:'
                    if not line.startswith(('A)', 'B)', 'C)', 'D)', 'Risposta corretta:')):
                        if current_question:
                            # Verifica che la domanda abbia esattamente 4 opzioni e una risposta corretta
                            if len(current_question['opzioni']) == 4 and corretta_opzione:
                                for opzione in current_question['opzioni']:
                                    if opzione['testo_opzione'].startswith(corretta_opzione):
                                        opzione['corretta'] = True
                                questionario.append(current_question)
                                domande_gia_generate.append(current_question)  # Aggiungi al registro delle domande
                                domande_generate += 1
                            else:
                                print(f"Domanda incompleta o senza risposta corretta: {current_question['testo_domanda']}")
                            if domande_generate >= num_domande:
                                break

                        # Inizia una nuova domanda
                        current_question = {"testo_domanda": line, "opzioni": []}
                        corretta_opzione = None

                    elif line.startswith(('A)', 'B)', 'C)', 'D)')) and current_question:
                        current_question['opzioni'].append({"testo_opzione": line, "corretta": False})
                    elif line.startswith('Risposta corretta:') and current_question:
                        corretta_opzione = line.split(':')[1].strip()

                # Verifica che la domanda corrente sia completa
                if current_question and domande_generate < num_domande:
                    if len(current_question['opzioni']) == 4 and corretta_opzione:
                        for opzione in current_question['opzioni']:
                            if opzione['testo_opzione'].startswith(corretta_opzione):
                                opzione['corretta'] = True
                        questionario.append(current_question)
                        domande_gia_generate.append(current_question)  # Aggiungi al registro delle domande
                        domande_generate += 1

            if domande_generate >= num_domande:
                break

    except Exception as e:
        print(f"Errore durante la generazione: {e}")
        return [], []

    return macro_aree[:5], questionario



def salva_macro_aree_e_questionario(lezione_id, macro_aree, questionario):
    conn = get_db_connection()
    if conn is None:
        print("Errore di connessione al database.")
        return False

    try:
        cursor = conn.cursor()

        # Salvataggio delle macro-aree nella tabella 'lezione_argomento'
        for nome_argomento, descrizione_argomento in macro_aree:
            cursor.execute("""
                INSERT INTO lezione_argomento (lezione_id, nome_argomento, descrizione_argomento)
                VALUES (%s, %s, %s)
            """, (lezione_id, nome_argomento, descrizione_argomento))

        # Creazione di un nuovo questionario per la lezione
        cursor.execute("""
            INSERT INTO questionario (lezione_id)
            VALUES (%s)
        """, (lezione_id,))
        questionario_id = cursor.lastrowid

        # Salvataggio delle domande e delle opzioni del questionario
        for domanda in questionario:
            # Inserisci la domanda
            cursor.execute("""
                INSERT INTO domanda (questionario_id, testo_domanda, corretta_opzione_id)
                VALUES (%s, %s, NULL)  # L'opzione corretta sarà aggiornata successivamente
            """, (questionario_id, domanda['testo_domanda']))
            domanda_id = cursor.lastrowid

            # Inserisci le opzioni della domanda
            for opzione in domanda['opzioni']:
                cursor.execute("""
                    INSERT INTO opzione (domanda_id, testo_opzione)
                    VALUES (%s, %s)
                """, (domanda_id, opzione['testo_opzione']))

                # Se questa è l'opzione corretta, aggiorna 'corretta_opzione_id' nella tabella 'domanda'
                if opzione['corretta']:
                    corretta_opzione_id = cursor.lastrowid
                    cursor.execute("""
                        UPDATE domanda
                        SET corretta_opzione_id = %s
                        WHERE domanda_id = %s
                    """, (corretta_opzione_id, domanda_id))

        # Commit delle operazioni sul database
        conn.commit()
        cursor.close()
        conn.close()
        return True

    except mysql.connector.Error as err:
        print(f"Errore durante l'inserimento nel database: {err}")
        conn.rollback()
        cursor.close()
        conn.close()
        return False


# Rotta per avviare una Lezione Programmata o Immediata
@app.route('/avvia_lezione', methods=['POST'])
def avvia_lezione():
    try:
        if 'user_id' in session and session['user_type'] == 'docente':
            corso_id = request.form.get('corso_id')

            conn = get_db_connection()
            if conn is None:
                # Se la connessione fallisce, mostra un messaggio di errore e reindirizza al login
                flash('Errore di connessione al database.')
                return redirect(url_for('login'))

            # Cursore per eseguire le query SQL
            cursor = conn.cursor(dictionary=True)

            docente_id = session['user_id']

            if 'type_start' in request.form and request.form['type_start'] == "Immediata":
                oggi = datetime.now().date()  # Ottiene la data odierna
                descrizione = request.form.get('descrizione')
                query = ("INSERT INTO lezione(corso_id, docente_id, data, descrizione, statoLezione) "
                         "VALUES (%s, %s, %s, %s, '1')")
                cursor.execute(query, (corso_id, docente_id, oggi, descrizione,))

                query = (
                    "SELECT lezione_id FROM lezione WHERE corso_id = %s AND docente_id = %s AND data = %s AND statoLezione = %s")
                cursor.execute(query, (corso_id, docente_id, oggi, '1',))

                # Ottiene i risultati della query (corsi)
                lezione_id = cursor.fetchone()['lezione_id']
            elif 'type_start' in request.form and request.form['type_start'] == "Programmata":
                lezione_id = request.form.get('lezione_id')

            # Verifica la presenza del file nel form
            if 'file' not in request.files or request.files['file'].filename == '':
                flash("Nessun file selezionato", 'error')
                return redirect(url_for('dashboard'))

            file = request.files['file']

            # Verifica estensione file
            if file and (file.filename.endswith('.pdf') or file.filename.endswith('.pptx')):
                filepath = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
                file.save(filepath)

                testo = ""

                # Estrazione del testo dal file
                if file.filename.endswith('.pdf'):
                    testo = estrai_testo_da_pdf(filepath)
                elif file.filename.endswith('.pptx'):
                    testo = estrai_testo_da_pptx(filepath)

                # Generazione delle macro-aree e del questionario in un unico blocco
                macro_aree, questionario = genera_macro_aree_e_questionario(testo)

                # Salvataggio nel database
                if salva_macro_aree_e_questionario(lezione_id, macro_aree, questionario):
                    flash("Lezione avviata con successo e questionario generato.", 'success')
                else:
                    flash("Errore durante il salvataggio nel database.", 'error')

                query = "UPDATE lezione SET statoLezione= '2' WHERE lezione_id = %s"
                cursor.execute(query, (lezione_id,))

                cursor.close()  # Chiude il cursore
                conn.close()  # Chiude la connessione al database

                return redirect(url_for('accedi_lezione', lezione_id=lezione_id, corso_id=corso_id))


            else:
                flash("Formato file non valido. Si accettano solo PDF o PPTX", 'error')
                return redirect(url_for('lezioni'))
        else:
            flash('Non sei autorizzato ad avviare una lezione.', 'error')
            return redirect(url_for('dashboard'))
    except Exception as e:
        print(f"Errore durante l'avvio della lezione: {e}")
        return "Internal Server Error", 500


@app.route('/accedi_lezione', methods=['GET', 'POST'])
def accedi_lezione():
    # Verifica se si tratta di una richiesta GET
    if request.method == 'GET':
        lezione_id = request.args.get('lezione_id')
        corso_id = request.args.get('corso_id')
    elif request.method == 'POST':
        # Recupera i dati dalla richiesta POST
        lezione_id = request.form.get('lezione_id')
        corso_id = request.form.get('corso_id')

    if 'user_id' in session:
        user_id = session['user_id']
        user_type = session['user_type']

        # Connessione al database
        conn = get_db_connection()
        if conn is None:
            return jsonify({"error": "Errore di connessione al database"}), 500

        cursor = conn.cursor(dictionary=True)

        # Assicurati di leggere il risultato della query
        query = "SELECT questionario_id FROM questionario WHERE lezione_id = %s"
        cursor.execute(query, (lezione_id,))
        questionario_result = cursor.fetchone()  # Modificato per assicurarsi che i risultati siano letti
        questionario_id = questionario_result['questionario_id'] if questionario_result else None
        cursor.fetchall()

        if user_type == 'studente':
            query = "SELECT * FROM Presente WHERE studente_id = %s AND lezione_id = %s"
            cursor.execute(query, (user_id, lezione_id,))
            result = cursor.fetchone()  # Modificato per assicurarsi che i risultati siano letti

            if result is None:
                query = "INSERT INTO Presente (studente_id, lezione_id) VALUES (%s, %s)"
                cursor.execute(query, (user_id, lezione_id,))
                conn.commit()

            cursor.close()
            conn.close()

            return render_template('lezioneStudente.html', user_id=user_id, user_type=user_type,
                                   corso_id=corso_id, lezione_id=lezione_id, questionario_id=questionario_id)
        else:
            cursor.close()
            conn.close()

            return render_template('lezioneDocente.html', user_id=user_id, user_type=user_type,
                                   corso_id=corso_id, lezione_id=lezione_id, questionario_id=questionario_id)

    else:
        return redirect(url_for('login'))


@app.route('/aggiorna_dati', methods=['POST'])
def aggiorna_dati():
    lezione_id = request.json.get('lezione_id')

    # Connessione al database
    conn = get_db_connection()
    if conn is None:
        return jsonify({"error": "Errore di connessione al database"}), 500

    cursor = conn.cursor(dictionary=True)

    # 1. Ottenere l'elenco degli studenti presenti alla lezione
    query_studenti = """
        SELECT s.nome, s.cognome 
        FROM Studente s 
        JOIN Presente p ON s.studente_id = p.studente_id 
        WHERE p.lezione_id = %s
    """
    cursor.execute(query_studenti, (lezione_id,))
    studenti_presenti = cursor.fetchall()

    # 2. Ottenere le domande e le opzioni della lezione
    query_domande = """
        SELECT d.domanda_id, d.testo_domanda, d.corretta_opzione_id
        FROM Domanda d
        JOIN Questionario q ON d.questionario_id = q.questionario_id
        WHERE q.lezione_id = %s
    """
    cursor.execute(query_domande, (lezione_id,))
    domande = cursor.fetchall()

    # Aggiungere le opzioni per ogni domanda
    for domanda in domande:
        query_opzioni = """
            SELECT o.testo_opzione, o.opzione_id
            FROM Opzione o
            WHERE o.domanda_id = %s
        """
        cursor.execute(query_opzioni, (domanda['domanda_id'],))
        domanda['opzioni'] = cursor.fetchall()

    # 3. Ottenere gli argomenti della lezione
    query_argomenti = """
        SELECT la.nome_argomento, la.descrizione_argomento
        FROM Lezione_Argomento la
        WHERE la.lezione_id = %s
    """
    cursor.execute(query_argomenti, (lezione_id,))
    argomenti = cursor.fetchall()

    # 4. Ottenere lo stato della lezione dalla tabella AttivitàLezione
    query_stato_lezione = """
        SELECT al.stato_lezione, al.modalità_lezione, al.fase_lezione, al.descrizione
        FROM AttivitàLezione al
        JOIN Lezione l ON al.attività_lezione_id = l.statoLezione
        WHERE l.lezione_id = %s
    """
    cursor.execute(query_stato_lezione, (lezione_id,))
    stato_lezione_info = cursor.fetchone()

    cursor.close()
    conn.close()

    return jsonify({
        'studenti_presenti': studenti_presenti,
        'domande': domande,
        'argomenti': argomenti,
        'stato_lezione': stato_lezione_info['stato_lezione'],
        'modalità_lezione': stato_lezione_info['modalità_lezione'],
        'fase_lezione': stato_lezione_info['fase_lezione'],
        'descrizione': stato_lezione_info['descrizione']
    })


@app.route('/modifica_fase_lezione', methods=['POST'])
def modifica_fase_lezione():
    data = request.get_json()
    lezione_id = data.get('lezione_id')
    questionario_id = data.get('questionario_id')

    try:
        conn = get_db_connection()
        cursor = conn.cursor()

        # 1. Recupera l'attuale attività_lezione_id della lezione
        query_attivita_corrente = """
            SELECT al.modalità_lezione, al.fase_lezione, al.attività_lezione_id
            FROM AttivitàLezione al
            JOIN Lezione l ON al.attività_lezione_id = l.statoLezione
            WHERE l.lezione_id = %s
        """
        cursor.execute(query_attivita_corrente, (lezione_id,))
        risultato_corrente = cursor.fetchone()

        if risultato_corrente is None:
            return jsonify({'success': False, 'error': 'Lezione non trovata.'}), 404

        modalità_corrente = risultato_corrente[0]
        fase_corrente = risultato_corrente[1]
        attività_lezione_id_corrente = risultato_corrente[2]

        # 2. Trova la prossima attività_lezione_id in base alla modalità e alla fase corrente
        query_prossima_fase = """
            SELECT attività_lezione_id
            FROM AttivitàLezione
            WHERE modalità_lezione = %s AND attività_lezione_id = %s + 1
            ORDER BY fase_lezione ASC
            LIMIT 1
        """
        cursor.execute(query_prossima_fase, (modalità_corrente, attività_lezione_id_corrente))
        nuova_fase_id = cursor.fetchone()

        if nuova_fase_id is None:
            return jsonify({'success': False, 'error': 'Non ci sono altre fasi disponibili per questa modalità.'}), 400

        nuova_fase_id = nuova_fase_id[0]

        # 3. Aggiorna lo stato della lezione nella tabella Lezione
        query_update_fase = """
            UPDATE Lezione
            SET statoLezione = %s
            WHERE lezione_id = %s
        """
        cursor.execute(query_update_fase, (nuova_fase_id, lezione_id))

        data_ora_avvio = datetime.now()

        if nuova_fase_id == 3:
            query_inizio_questionario = """
                        UPDATE Questionario
                        SET avvio_primaFase = %s
                        WHERE lezione_id = %s AND questionario_id = %s 
                    """
            cursor.execute(query_inizio_questionario, (data_ora_avvio, lezione_id, questionario_id))

        if nuova_fase_id == 5:
            query_inizio_questionario = """
                        UPDATE Questionario
                        SET avvio_secondaFase = %s
                        WHERE lezione_id = %s AND questionario_id = %s 
                    """
            cursor.execute(query_inizio_questionario, (data_ora_avvio, lezione_id, questionario_id))

        conn.commit()
        cursor.close()
        conn.close()
        return jsonify({'success': True})
    except Exception as e:
        print(f"Errore: {e}")
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/termina_lezione', methods=['POST'])
def termina_lezione():
    data = request.get_json()
    lezione_id = data.get('lezione_id')

    try:
        conn = get_db_connection()
        cursor = conn.cursor()
        # Imposta lo stato della lezione come "terminata"
        query = "UPDATE Lezione SET statoLezione = (SELECT attività_lezione_id FROM AttivitàLezione WHERE stato_lezione = 'terminata') WHERE lezione_id = %s"
        cursor.execute(query, (lezione_id,))
        conn.commit()
        cursor.close()
        conn.close()
        return jsonify({'success': True})
    except Exception as e:
        print(f"Errore: {e}")
        return jsonify({'success': False})


# Rotta per avviare una Lezione Programmata o Immediata
@app.route('/programma_lezione', methods=['POST'])
def programma_lezione():
    try:
        if 'user_id' in session and session['user_type'] == 'docente':

            conn = get_db_connection()
            if conn is None:
                # Se la connessione fallisce, mostra un messaggio di errore e reindirizza al login
                flash('Errore di connessione al database.')
                return redirect(url_for('login'))

            # Cursore per eseguire le query SQL
            cursor = conn.cursor(dictionary=True)

            docente_id = session['user_id']
            corso_id = request.form.get('corso_id')
            data = request.form.get('date')
            descrizione = request.form.get('descrizione')
            query = ("INSERT INTO lezione(corso_id, docente_id, data, descrizione, statoLezione) "
                     "VALUES (%s, %s, %s, %s, '1')")
            cursor.execute(query, (corso_id, docente_id, data, descrizione,))

            cursor.close()  # Chiude il cursore
            conn.close()  # Chiude la connessione al database

            return redirect(url_for('indexCorso', corso_id=corso_id))
        else:
            flash('Non sei autorizzato a programmare una lezione.', 'error')
            return redirect(url_for('dashboard'))
    except Exception as e:
        print(f"Errore durante programmazione della lezione: {e}")
        return "Internal Server Error", 500


# Funzione che gestisce la connessione al database MySQL
def get_db_connection():
    try:
        # Crea la connessione con il database
        conn = mysql.connector.connect(
            host="localhost",  # Indirizzo del server MySQL
            user="root",  # Nome utente del database MySQL
            password="",  # Password per l'utente del database
            database="FlipInClassIA_DB"  # Nome del database da utilizzare
        )
        return conn
    except mysql.connector.Error as err:
        # Se c'è un errore nella connessione, lo stampa e restituisce None
        print(f"Errore di connessione al database: {err}")
        return None


# Rotta per la homepage dell'app
@app.route('/')
def index():
    # Restituisce la pagina principale ('index.html') come risposta
    return render_template('index.html')


# Rotta per la pagina di login
@app.route('/login')
def login():
    # Restituisce la pagina di login ('login.html')
    return render_template('login.html')


# Rotta per la pagina di registrazione
@app.route('/registrati')
def registrati():
    # Restituisce la pagina di registrazione ('registrati.html')
    return render_template('registrati.html')


# Rotta per la dashboard, visibile solo agli utenti loggati
@app.route('/dashboard')
def dashboard():
    # Verifica se l'utente è loggato controllando se c'è un 'user_id' nella sessione
    if 'user_id' in session:
        user_id = session['user_id']  # Recupera l'ID utente dalla sessione
        user_type = session['user_type']  # Recupera il tipo di utente (docente o studente)

        # Connessione al database
        conn = get_db_connection()
        if conn is None:
            # Se la connessione fallisce, mostra un messaggio di errore e reindirizza al login
            flash('Errore di connessione al database.')
            return redirect(url_for('login'))

        # Cursore per eseguire le query SQL
        cursor = conn.cursor(dictionary=True)

        # Se l'utente è un docente, recupera i corsi associati a quel docente
        if user_type == 'docente':
            query = (
                "SELECT C.corso_id, C.nome AS corso_nome, C.descrizione AS corso_descrizione, C.image_path AS corso_immagine "
                "FROM Corso AS C WHERE C.corso_id IN (SELECT corso_id FROM Lavora WHERE docente_id = %s)")
            cursor.execute(query, (user_id,))  # Esegue la query passando l'ID del docente
        else:
            # Se l'utente è uno studente, recupera i corsi a cui lo studente partecipa
            query = (
                "SELECT C.corso_id, C.nome AS corso_nome, C.descrizione AS corso_descrizione, C.image_path AS corso_immagine "
                "FROM Corso AS C "
                "JOIN Partecipa AS P ON C.corso_id = P.corso_id "
                "WHERE P.studente_id = %s")
            cursor.execute(query, (user_id,))  # Esegue la query con l'ID dello studente

        # Ottiene i risultati della query (corsi)
        corsi = cursor.fetchall()
        cursor.close()  # Chiude il cursore
        conn.close()  # Chiude la connessione al database

        # Restituisce la pagina della dashboard con i corsi dell'utente
        return render_template('dashboard.html', user_id=user_id, user_type=user_type, corsi=corsi)

    else:
        # Se l'utente non è loggato, mostra un messaggio e reindirizza alla pagina di login
        flash('Devi essere loggato per accedere alla dashboard.')
        return redirect(url_for('login'))


# Funzione per generare un'immagine personalizzata per un corso
def generate_course_image(course_name, corso_id):
    # Crea una nuova immagine di 200x200 pixel con un colore di sfondo casuale
    img = Image.new('RGB', (200, 200), color=(random.randint(0, 255), random.randint(0, 255), random.randint(0, 255)))

    # Usa le prime due lettere del nome del corso (in maiuscolo) come iniziali
    initials = (course_name[:2] if len(course_name) >= 2 else course_name).upper()

    # Crea un oggetto di disegno e tenta di caricare il font 'arial.ttf'
    draw = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("arial.ttf", 100)
    except IOError:
        # Se il font non è disponibile, usa un font di default
        font = ImageFont.load_default()

    # Calcola la posizione per centrare il testo nell'immagine usando textbbox
    bbox = draw.textbbox((0, 0), initials, font=font)
    text_width = bbox[2] - bbox[0]
    text_height = bbox[3] - bbox[1]
    position = ((img.width - text_width) // 2, (img.height - text_height) // 2)

    # Disegna le iniziali al centro dell'immagine
    draw.text(position, initials, fill=(255, 255, 255), font=font)

    # Salva l'immagine con un nome basato sul corso_id
    image_filename = f"{corso_id}.png"
    image_path = os.path.join(IMAGES_PATH, image_filename)
    img.save(image_path)  # Salva l'immagine

    return image_filename  # Restituisce il nome del file immagine da inserire nel database


# Rotta per creare un nuovo corso, accessibile solo ai docenti
@app.route('/crea_corso', methods=['POST'])
def crea_corso():
    # Controlla se l'utente è loggato e se è un docente
    if 'user_id' in session and session['user_type'] == 'docente':
        course_name = request.form['courseName']  # Ottiene il nome del corso dal form HTML
        course_description = request.form['courseDescription']  # Ottiene la descrizione del corso
        docente_id = session['user_id']  # Recupera l'ID del docente dalla sessione

        # Genera un ID univoco per il corso (un codice esadecimale di 6 caratteri)
        corso_id = ''.join(random.choices('0123456789abcdef', k=6))

        # Genera un'immagine personalizzata per il corso
        image_filename = generate_course_image(course_name, corso_id)

        # Connessione al database
        conn = get_db_connection()
        cursor = conn.cursor()

        # Inserisce il nuovo corso nel database con il percorso dell'immagine
        query = "INSERT INTO Corso (corso_id, nome, descrizione, image_path) VALUES (%s, %s, %s, %s)"
        cursor.execute(query, (corso_id, course_name, course_description, image_filename))

        # Inserisce un record che associa il docente creatore come proprietario del corso
        query = "INSERT INTO Lavora (docente_id, corso_id, proprietario) VALUES (%s, %s, 1)"
        cursor.execute(query, (docente_id, corso_id))

        # Esegue le modifiche e chiude il cursore e la connessione
        conn.commit()
        cursor.close()
        conn.close()

    else:
        # Se non è un docente loggato, mostra un messaggio di errore
        flash('Non sei autorizzato a creare un corso.')

    # Reindirizza alla dashboard dopo la creazione del corso
    return redirect(url_for('dashboard'))


# Rotta per partecipare a un corso, accessibile solo agli studenti
@app.route('/partecipa_corso', methods=['POST'])
def partecipa_corso():
    # Verifica se l'utente è loggato e se è uno studente
    if 'user_id' in session and session['user_type'] == 'studente':
        course_code = request.form['courseCode']  # Ottiene il codice del corso dal form
        studente_id = session['user_id']  # Recupera l'ID dello studente dalla sessione

        # Connessione al database
        conn = get_db_connection()
        cursor = conn.cursor()

        # Inserisce il record di partecipazione dello studente al corso
        query = "INSERT INTO Partecipa (studente_id, corso_id) VALUES (%s, %s)"
        cursor.execute(query, (studente_id, course_code))

        # Esegue le modifiche e chiude il cursore e la connessione
        conn.commit()
        cursor.close()
        conn.close()

    else:
        # Se non è uno studente loggato, mostra un messaggio di errore
        flash('Non sei autorizzato a partecipare a un corso.')

    # Reindirizza alla dashboard
    return redirect(url_for('dashboard'))


# Funzione per ottenere il docente proprietario di un corso
def ottieniProfessorePresidente(corso_id):
    # Connessione al database
    conn = get_db_connection()
    cursor = conn.cursor(dictionary=True)

    # Esegue una query per ottenere il nome, cognome e immagine del docente che è proprietario del corso
    query = ("SELECT D.nome, D.cognome, D.image_path FROM Lavora AS L, Docente AS D "
             "WHERE L.corso_id = %s AND L.proprietario = 1 AND D.docente_id = L.docente_id")
    cursor.execute(query, (corso_id,))
    docente_presidente = cursor.fetchone()  # Ottiene il risultato della query

    return docente_presidente  # Restituisce le informazioni del docente


# Funzione per ottenere il nome di un corso dato il suo ID
def ottieniNomeCorso(corso_id):
    # Connessione al database
    conn = get_db_connection()
    cursor = conn.cursor(dictionary=True)

    # Esegue una query per ottenere il nome del corso
    query = ("SELECT C.nome FROM Corso AS C WHERE C.corso_id = %s")
    cursor.execute(query, (corso_id,))
    nome_corso = cursor.fetchone()  # Ottiene il nome del corso

    if nome_corso:
        nome_corso = nome_corso['nome']  # Estrae il nome dal risultato
    else:
        nome_corso = "Corso non trovato"  # Messaggio di errore se il corso non è stato trovato

    return nome_corso  # Restituisce il nome del corso


# Rotta per la pagina dettagliata di un corso specifico
@app.route('/indexCorso')
def indexCorso():
    # Controlla se l'utente è loggato
    if 'user_id' in session:
        user_id = session['user_id']  # Recupera l'ID dell'utente dalla sessione
        user_type = session['user_type']  # Recupera il tipo di utente (studente o docente)

        # Recupera il corso_id dalla query string (parametro URL)
        corso_id = request.args.get('corso_id')

        oggi = datetime.now().date()  # Ottiene la data odierna

        # Connessione al database
        conn = get_db_connection()
        if conn is None:
            return jsonify({
                "error": "Errore di connessione al database"}), 500  # Messaggio di errore in caso di fallimento della connessione

        cursor = conn.cursor(dictionary=True)

        query = ("SELECT COUNT(LE.lezione_id) AS NumLezioni "
                 "FROM Lezione AS LE "
                 "JOIN Docente AS D ON LE.docente_id = D.docente_id "
                 "JOIN Corso AS C ON C.corso_id = LE.corso_id "
                 "WHERE C.corso_id = %s AND LE.data > %s")
        cursor.execute(query, (corso_id, oggi))
        NumLezioni = cursor.fetchone()['NumLezioni']  # Recupera tutte le lezioni del corso

        # Query per ottenere il calendario delle lezioni del corso
        query = ("SELECT D.nome, D.cognome, LE.data, LE.descrizione, LE.lezione_id "
                 "FROM Lezione AS LE "
                 "JOIN Docente AS D ON LE.docente_id = D.docente_id "
                 "JOIN Corso AS C ON C.corso_id = LE.corso_id "
                 "WHERE C.corso_id = %s "
                 "ORDER BY LE.data DESC")
        cursor.execute(query, (corso_id,))
        calendario_lezioni = cursor.fetchall()  # Recupera tutte le lezioni del corso

        # Query per ottenere gli argomenti trattati in ogni lezione del corso
        query = ("SELECT ARG.nome_argomento, ARG.lezione_id "
                 "FROM Lezione_argomento AS ARG, Corso AS C, Lezione AS L "
                 "WHERE C.corso_id = %s "
                 "AND C.corso_id = L.corso_id AND L.lezione_id = ARG.lezione_id ")
        cursor.execute(query, (corso_id,))
        argomenti_lezioni = cursor.fetchall()  # Recupera gli argomenti delle lezioni

        docente_presidente = ottieniProfessorePresidente(corso_id)  # Ottiene il docente proprietario del corso
        nome_corso = ottieniNomeCorso(corso_id)  # Ottiene il nome del corso

        # Chiude i cursori e la connessione
        cursor.close()
        conn.close()

        # Restituisce la pagina del corso ('indexCorso.html') con tutte le informazioni necessarie
        return render_template('indexCorso.html', user_id=user_id, user_type=user_type, NumLezioni=NumLezioni,
                               nome_corso=nome_corso, corso_id=corso_id, calendario_lezioni=calendario_lezioni,
                               oggi=oggi, argomenti_lezioni=argomenti_lezioni, docente_presidente=docente_presidente)

    else:
        # Se l'utente non è loggato, mostra un messaggio e reindirizza al login
        flash('Devi essere loggato per accedere alla dashboard.')
    return redirect(url_for('login'))


# Rotta per visualizzare la lista degli studenti che partecipano a un corso
@app.route('/studentiPartecipanti')
def studentiPartecipanti():
    # Controlla se l'utente è loggato
    if 'user_id' in session:
        user_id = session['user_id']  # Recupera l'ID utente dalla sessione
        user_type = session['user_type']  # Recupera il tipo di utente (studente o docente)

        # Recupera il corso_id passato come parametro dalla query string
        corso_id = request.args.get('corso_id')

        # Connessione al database
        conn = get_db_connection()
        if conn is None:
            # In caso di errore nella connessione, restituisce un messaggio di errore come JSON
            return jsonify({"error": "Errore di connessione al database"}), 500

        cursor = conn.cursor(dictionary=True)

        # Esegue una query per ottenere le informazioni sugli studenti iscritti al corso
        query = ("SELECT S.nome, S.cognome, S.email, S.image_path "
                 "FROM Studente AS S, Partecipa AS P "
                 "WHERE P.corso_id = %s AND S.studente_id = P.studente_id")
        cursor.execute(query, (corso_id,))
        studenti_corso = cursor.fetchall()  # Recupera la lista di studenti

        docente_presidente = ottieniProfessorePresidente(corso_id)  # Ottiene il docente responsabile del corso
        nome_corso = ottieniNomeCorso(corso_id)  # Ottiene il nome del corso

        # Chiude il cursore e la connessione
        cursor.close()
        conn.close()

        # Restituisce la pagina con l'elenco degli studenti partecipanti ('studenti_partecipanti.html')
        return render_template('studenti_partecipanti.html', user_id=user_id, user_type=user_type,
                               corso_id=corso_id, studenti_corso=studenti_corso,
                               docente_presidente=docente_presidente, nome_corso=nome_corso)

    else:
        # Se l'utente non è loggato, mostra un messaggio di errore e reindirizza alla pagina di login
        flash('Devi essere loggato per accedere alla dashboard.')
    return redirect(url_for('login'))


# Rotta per visualizzare la lista dei professori che partecipano a un corso
@app.route('/professoriPartecipanti')
def professoriPartecipanti():
    # Controlla se l'utente è loggato
    if 'user_id' in session:
        user_id = session['user_id']  # Recupera l'ID utente dalla sessione
        user_type = session['user_type']  # Recupera il tipo di utente (studente o docente)

        # Recupera il corso_id passato come parametro dalla query string
        corso_id = request.args.get('corso_id')

        # Connessione al database
        conn = get_db_connection()
        if conn is None:
            # In caso di errore nella connessione, restituisce un messaggio di errore come JSON
            return jsonify({"error": "Errore di connessione al database"}), 500

        cursor = conn.cursor(dictionary=True)

        # Esegue una query per ottenere le informazioni sui professori associati al corso
        query = ("SELECT D.nome, D.cognome, D.email, D.image_path, L.proprietario "
                 "FROM Docente AS D, Lavora AS L "
                 "WHERE L.corso_id = %s AND D.docente_id = L.docente_id")
        cursor.execute(query, (corso_id,))
        docente_corso = cursor.fetchall()  # Recupera la lista di docenti

        docente_presidente = ottieniProfessorePresidente(corso_id)  # Ottiene il docente responsabile del corso
        nome_corso = ottieniNomeCorso(corso_id)  # Ottiene il nome del corso

        # Chiude il cursore e la connessione
        cursor.close()
        conn.close()

        # Restituisce la pagina con l'elenco dei professori partecipanti ('professori_partecipanti.html')
        return render_template('professori_partecipanti.html', user_id=user_id, user_type=user_type,
                               corso_id=corso_id, docente_corso=docente_corso,
                               docente_presidente=docente_presidente, nome_corso=nome_corso)

    else:
        # Se l'utente non è loggato, mostra un messaggio di errore e reindirizza alla pagina di login
        flash('Devi essere loggato per accedere alla dashboard.')
    return redirect(url_for('login'))


# Rotta per visualizzare e/o partecipare alle lezioni giornaliere del corso
@app.route('/lezioni')
def lezioni():
    # Controlla se l'utente è loggato
    if 'user_id' in session:
        user_id = session['user_id']  # Recupera l'ID utente dalla sessione
        user_type = session['user_type']  # Recupera il tipo di utente (studente o docente)

        # Recupera il corso_id passato come parametro dalla query string
        corso_id = request.args.get('corso_id')

        # Connessione al database
        conn = get_db_connection()
        if conn is None:
            # In caso di errore nella connessione, restituisce un messaggio di errore come JSON
            return jsonify({"error": "Errore di connessione al database"}), 500

        cursor = conn.cursor(dictionary=True)

        # Ottiene la data corrente per il confronto
        oggi = datetime.now().date()

        if user_type == "docente":
            query = ("SELECT D.nome, D.cognome, LE.descrizione, LE.statoLezione, LE.lezione_id "
                     "FROM Lezione AS LE "
                     "JOIN Docente AS D ON LE.docente_id = D.docente_id "
                     "JOIN Corso AS C ON C.corso_id = LE.corso_id "
                     "WHERE C.corso_id = %s AND LE.data = %s AND LE.docente_id = %s ")
            cursor.execute(query, (corso_id, oggi, user_id))
        else:
            query = ("SELECT D.nome, D.cognome, LE.descrizione, LE.statoLezione, LE.lezione_id "
                     "FROM Lezione AS LE "
                     "JOIN Docente AS D ON LE.docente_id = D.docente_id "
                     "JOIN Corso AS C ON C.corso_id = LE.corso_id "
                     "WHERE C.corso_id = %s AND LE.data = %s")
            cursor.execute(query, (corso_id, oggi,))

        lezioni = cursor.fetchall()  # Recupera tutte le lezioni del corso

        docente_presidente = ottieniProfessorePresidente(corso_id)  # Ottiene il docente responsabile del corso
        nome_corso = ottieniNomeCorso(corso_id)  # Ottiene il nome del corso

        # Chiude il cursore e la connessione
        cursor.close()
        conn.close()

        # Restituisce la pagina con l'elenco dei professori partecipanti ('professori_partecipanti.html')
        return render_template('indexLezioni.html', user_id=user_id, user_type=user_type,
                               corso_id=corso_id, lezioni=lezioni,
                               docente_presidente=docente_presidente, nome_corso=nome_corso, oggi=oggi)

    else:
        # Se l'utente non è loggato, mostra un messaggio di errore e reindirizza alla pagina di login
        flash('Devi essere loggato per accedere alla dashboard.')
    return redirect(url_for('login'))


# Rotta per ottenere le informazioni delle lezioni in formato JSON (API)
@app.route('/api/lezioni', methods=['GET'])
def get_lezioni():
    # Ottiene il corso_id dai parametri della query string
    corso_id = request.args.get('corso_id')
    if not corso_id:
        # Se manca il parametro 'corso_id', restituisce un errore
        return jsonify({"error": "Parametro 'corso_id' mancante"}), 400

    # Connessione al database
    conn = get_db_connection()
    if conn is None:
        # Se c'è un errore nella connessione, restituisce un errore
        return jsonify({"error": "Errore di connessione al database"}), 500

    cursor = conn.cursor(dictionary=True)

    # Query per recuperare tutte le lezioni di un determinato corso
    query = "SELECT * FROM lezione WHERE corso_id = %s"
    cursor.execute(query, (corso_id,))

    # Ottiene la data corrente per il confronto
    oggi = datetime.now().date()

    lezioni = cursor.fetchall()  # Recupera tutte le lezioni del corso
    result = []

    # Itera attraverso le lezioni e determina il colore in base alla data
    for lezione in lezioni:
        try:
            data_lezione = lezione['data']  # Ottiene la data della lezione
            if isinstance(data_lezione, str):
                # Se la data è una stringa, la converte in oggetto `datetime.date`
                data_lezione = datetime.strptime(data_lezione, '%Y-%m-%d').date()
            elif isinstance(data_lezione, datetime):
                data_lezione = data_lezione.date()

            # Determina il colore dell'evento in base alla data
            if data_lezione < oggi:
                colore = 'green'  # Lezioni passate
                calendar = 'Lezione passata'
            elif data_lezione == oggi:
                colore = 'orange'  # Lezioni odierne
                calendar = 'Lezione odierna'
            else:
                colore = 'blue'  # Lezioni future
                calendar = 'Lezione programmata'

            # Aggiunge il risultato in formato JSON-friendly
            result.append({
                'descrizione': f"{lezione['descrizione']}",
                'calendar': calendar,
                'color': colore,
                'date': data_lezione.strftime('%Y-%m-%d')
            })

        except Exception as e:
            # In caso di errore durante il processo di data, stampa l'errore
            print(f"Errore nella gestione della data: {e}")

    # Chiude il cursore e la connessione
    cursor.close()
    conn.close()

    # Restituisce il risultato in formato JSON
    return jsonify(result)


# Rotta per gestire il form di login
@app.route('/loginForm', methods=['GET', 'POST'])
def loginForm():
    if request.method == 'POST':
        email = request.form['email']  # Ottiene l'email dal form
        password = request.form['password']  # Ottiene la password dal form

        # Connessione al database
        conn = get_db_connection()
        if conn is None:
            flash('Errore di connessione al database.')
            return redirect(url_for('login'))

        cursor = conn.cursor(dictionary=True)

        # Controlla se l'utente è un docente
        query = "SELECT * FROM Docente WHERE email = %s"
        cursor.execute(query, (email,))
        user = cursor.fetchone()

        # Se l'utente non è un docente, controlla se è uno studente
        if user is None:
            query = "SELECT * FROM Studente WHERE email = %s"
            cursor.execute(query, (email,))
            user = cursor.fetchone()
            user_type = 'studente'
        else:
            user_type = 'docente'

        conn.close()

        # Se l'utente esiste e la password è corretta
        if user and bcrypt.checkpw(password.encode('utf-8'), user['password'].encode('utf-8')):
            # Imposta i dettagli dell'utente nella sessione
            session['user_id'] = user['docente_id'] if user_type == 'docente' else user['studente_id']
            session['user_type'] = user_type
            session['user_name'] = user['nome'] + " " + user['cognome']
            session['image'] = user['image_path']

            # Reindirizza alla dashboard
            return redirect(url_for('dashboard'))
        else:
            flash('Email o password non corretti. Riprova.')

    return render_template('login.html')


# Rotta per il logout dell'utente
@app.route('/logout')
def logout():
    # Cancella tutte le informazioni della sessione
    session.clear()
    return redirect(url_for('login'))


# Rotta per la visualizzazione del profilo utente
@app.route('/profilo')
def profilo():
    # Verifica se l'utente è loggato
    if 'user_id' in session:
        user_id = session['user_id']  # Ottiene l'ID dell'utente
        user_type = session['user_type']  # Ottiene il tipo di utente (docente o studente)

        # Connessione al database
        conn = get_db_connection()
        if conn is None:
            flash('Errore di connessione al database.')
            return redirect(url_for('registrati'))

        cursor = conn.cursor(dictionary=True)

        # Query per ottenere i dettagli dell'utente
        if user_type == 'docente':
            query = "SELECT * FROM Docente WHERE docente_id = %s"
        else:
            query = "SELECT * FROM Studente WHERE studente_id = %s"
        cursor.execute(query, (user_id,))
        profilo_utente = cursor.fetchone()

        cursor.close()  # Chiude il cursore
        conn.close()  # Chiude la connessione al database

        # Restituisce la pagina del profilo ('profilo.html')
        return render_template('profilo.html', user_id=user_id, user_type=user_type, profilo_utente=profilo_utente)

    else:
        # Se l'utente non è loggato, mostra un messaggio e reindirizza al login
        flash('Devi essere loggato per accedere alla dashboard.')
        return redirect(url_for('login'))


def generate_user_image(name, cognome, user_id):
    # Crea un'immagine con uno sfondo colorato casuale
    img = Image.new('RGB', (200, 200), color=(random.randint(0, 255), random.randint(0, 255), random.randint(0, 255)))

    # Usa le iniziali del nome e cognome
    initials = (name[0] + cognome[0]).upper()

    draw = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("arial.ttf", 100)
    except IOError:
        font = ImageFont.load_default()

    # Calcola la posizione per centrare il testo
    bbox = draw.textbbox((0, 0), initials, font=font)
    text_width = bbox[2] - bbox[0]
    text_height = bbox[3] - bbox[1]
    position = ((img.width - text_width) // 2, (img.height - text_height) // 2)

    # Disegna le iniziali sull'immagine
    draw.text(position, initials, fill=(255, 255, 255), font=font)

    # Salva l'immagine con il user_id come nome file
    image_filename = f"{user_id}.png"
    image_path = os.path.join(USER_IMAGES_PATH, image_filename)
    img.save(image_path)

    return image_filename  # Restituisci il nome del file per il database


# Rotta per la gestione del form di registrazione
@app.route('/registratiForm', methods=['GET', 'POST'])
def registratiForm():
    if request.method == 'POST':
        name = request.form['name']
        cognome = request.form['cognome']
        email = request.form['email']
        password = request.form['password']
        user_type = 'docente' if 'userType' in request.form else 'studente'

        # Connessione al database
        conn = get_db_connection()

        cursor = conn.cursor(dictionary=True)

        try:
            # Verifica se l'email esiste già per un docente
            query = "SELECT * FROM Docente WHERE email = %s"
            cursor.execute(query, (email,))
            existing_user = cursor.fetchone()  # Leggi un singolo risultato
            cursor.fetchall()  # This ensures that all results are fetched

            # Se non esiste tra i docenti, verifica se esiste tra gli studenti
            if not existing_user:
                query = "SELECT * FROM Studente WHERE email = %s"
                cursor.execute(query, (email,))
                existing_user = cursor.fetchone()  # Leggi un singolo risultato
                cursor.fetchall()  # Fetch remaining results

            # Se l'email esiste già, mostra un messaggio di errore
            if existing_user:
                flash('L\'email esiste già. Riprova!')
                return redirect(url_for('registrati'))
            # Genera un ID univoco per l'utente
            user_id = str(uuid.uuid4())
            # Cifra la password
            hashed_password = bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt())

            # Controlla se un file immagine è stato caricato
            image = request.files.get('profilePicture')
            if image and image.filename != '':
                # Se è stata caricata un'immagine, salvala
                ext = os.path.splitext(image.filename)[1]  # Estrai l'estensione del file
                image_filename = f"{user_id}{ext}"
                image_path = os.path.join(USER_IMAGES_PATH, image_filename)
                image.save(image_path)  # Salva l'immagine nel percorso specificato
            else:
                # Se non è stata caricata un'immagine, genera una immagine predefinita
                image_filename = generate_user_image(name, cognome, user_id)

            # Inserisci i dati nella tabella appropriata
            if user_type == 'docente':
                query = "INSERT INTO Docente (docente_id, nome, cognome, email, password, image_path) VALUES (%s, %s, %s, %s, %s, %s)"
                cursor.execute(query, (user_id, name, cognome, email, hashed_password.decode('utf-8'), image_filename))
            else:
                query = "INSERT INTO Studente (studente_id, nome, cognome, email, password, image_path) VALUES (%s, %s, %s, %s, %s, %s)"
                cursor.execute(query, (user_id, name, cognome, email, hashed_password.decode('utf-8'), image_filename))

            conn.commit()

        except mysql.connector.Error as err:
            flash(f'Errore durante la registrazione: {err}')
            return redirect(url_for('registrati'))

        finally:
            # Chiudi il cursore e la connessione
            cursor.close()
            conn.close()

        # Imposta i dettagli dell'utente nella sessione
        session['user_id'] = user_id
        session['user_type'] = user_type
        session['user_name'] = name + " " + cognome
        session['image'] = image_filename

        flash('Registrazione avvenuta con successo!')
        return redirect(url_for('dashboard'))

    return redirect(url_for('registrati'))


# Rotta per la modifica del profilo utente
@app.route('/modificaProfilo', methods=['GET', 'POST'])
def modificaProfilo():
    if request.method == 'POST':
        user_type = request.form.get('user_type')  # Ottiene il tipo di utente
        user_id = request.form.get('user_id')  # Ottiene l'ID dell'utente

        # Connessione al database
        conn = get_db_connection()
        if conn is None:
            flash('Errore di connessione al database.')
            return redirect(url_for('registrati'))

        cursor = conn.cursor(dictionary=True)

        # Ottiene le informazioni del profilo utente
        if user_type == 'docente':
            query = "SELECT * FROM Docente WHERE docente_id = %s"
        else:
            query = "SELECT * FROM Studente WHERE studente_id = %s"
        cursor.execute(query, (user_id,))
        profilo_utente = cursor.fetchone()

        # Recupera i nuovi dati o mantiene quelli attuali se non sono stati modificati
        name = request.form.get('name') or profilo_utente['nome']
        cognome = request.form.get('cognome') or profilo_utente['cognome']
        email = request.form.get('email') or profilo_utente['email']

        # Verifica se l'email è già presente nel database
        try:
            query = "SELECT * FROM Docente WHERE email = %s AND docente_id != %s"
            cursor.execute(query, (email, user_id))
            existing_user = cursor.fetchone()

            if not existing_user:
                query = "SELECT * FROM Studente WHERE email = %s AND studente_id != %s"
                cursor.execute(query, (email, user_id))
                existing_user = cursor.fetchone()

            if existing_user:
                flash('L\'email esiste già. Riprova!')
                return redirect(url_for('profilo'))

            # Verifica se è stata caricata una nuova immagine
            image = request.files.get('profilePicture')
            if image and image.filename != '':
                ext = os.path.splitext(image.filename)[1]  # Ottiene l'estensione del file immagine
                image_filename = f"{user_id}{ext}"
                image_path = os.path.join(USER_IMAGES_PATH, image_filename)
                image.save(image_path)  # Salva l'immagine nel percorso specificato
            else:
                # Mantiene l'immagine attuale se non viene caricata una nuova immagine
                image_filename = profilo_utente["image_path"]

            # Aggiorna il profilo utente nel database
            if user_type == 'docente':
                query = "UPDATE Docente SET nome = %s, cognome = %s, email = %s, image_path = %s WHERE docente_id = %s"
                cursor.execute(query, (name, cognome, email, image_filename, user_id))
            else:
                query = "UPDATE Studente SET nome = %s, cognome = %s, email = %s, image_path = %s WHERE studente_id = %s"
                cursor.execute(query, (name, cognome, email, image_filename, user_id))

            conn.commit()  # Salva le modifiche nel database

        except mysql.connector.Error as err:
            flash(f'Errore durante la modifica: {err}')
            return redirect(url_for('registrati'))

        finally:
            cursor.close()  # Chiude il cursore
            conn.close()  # Chiude la connessione al database

        # Aggiorna i dettagli dell'utente nella sessione
        session['user_id'] = user_id
        session['user_type'] = user_type
        session['user_name'] = name + " " + cognome
        session['image'] = image_filename

        # Reindirizza alla pagina del profilo
        return redirect(url_for('profilo'))

    return redirect(url_for('login'))


# Avvio dell'applicazione Flask
if __name__ == '__main__':
    app.run(debug=True)
